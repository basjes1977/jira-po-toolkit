import sys
import time
"""
Jira Sprint PowerPoint Generator
--------------------------------
This script fetches Jira sprint data and generates a PowerPoint presentation using a template.
It groups issues by label, displays issue details (with assignee avatars), and includes summary and upcoming slides.
All text is placed using template placeholders for consistent formatting.
"""

def get_upcoming_sprint_id():
    url = f"{JIRA_URL}/rest/agile/1.0/board/{BOARD_ID}/sprint?state=future"
    resp = requests.get(url, auth=(JIRA_EMAIL, JIRA_API_TOKEN))
    resp.raise_for_status()
    sprints = resp.json().get("values", [])
    if not sprints:
        return None
    return sprints[0]["id"]
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from collections import defaultdict
import os

# Load Jira credentials from .jira_environment
from pathlib import Path
from dotenv import dotenv_values

def load_jira_env():
    """
    Load Jira credentials from a .jira_environment file in the parent directory.
    Returns a dict with environment variables.
    """
    env_path = Path(__file__).parent / ".jira_environment"
    if env_path.exists():
        # Parse the file manually since it's not a standard .env
        env = {}
        with open(env_path) as f:
            for line in f:
                line = line.strip()
                if line.startswith("export "):
                    line = line[len("export "):]
                    if "=" in line:
                        k, v = line.split("=", 1)
                        env[k.strip()] = v.strip().strip('"')
        return env
    return {}

JIRA_ENV = load_jira_env()
JIRA_URL = JIRA_ENV.get("JT_JIRA_URL", "https://equinixjira.atlassian.net/").rstrip("/")
JIRA_EMAIL = JIRA_ENV.get("JT_JIRA_USERNAME")
JIRA_API_TOKEN = JIRA_ENV.get("JT_JIRA_PASSWORD")
BOARD_ID = JIRA_ENV.get("JT_JIRA_BOARD")

def get_current_sprint_id():
    """
    Get the ID of the current active sprint from Jira.
    """
    url = f"{JIRA_URL}/rest/agile/1.0/board/{BOARD_ID}/sprint?state=active"
    resp = requests.get(url, auth=(JIRA_EMAIL, JIRA_API_TOKEN))
    resp.raise_for_status()
    sprints = resp.json().get("values", [])
    if not sprints:
        raise Exception("No active sprint found.")
    return sprints[0]["id"]

def get_issues(sprint_id):
    """
    Fetch all issues for a given sprint ID from Jira.
    Returns a list of issue dicts.
    """
    url = f"{JIRA_URL}/rest/agile/1.0/sprint/{sprint_id}/issue"
    issues = []
    start_at = 0
    while True:
        params = {"startAt": start_at, "maxResults": 50}
        resp = requests.get(url, params=params, auth=(JIRA_EMAIL, JIRA_API_TOKEN))
        resp.raise_for_status()
        data = resp.json()
        issues.extend(data["issues"])
        if start_at + 50 >= data["total"]:
            break
        start_at += 50
    return issues

def get_sprint_name(sprint_id):
    """
    Fetch the name of a sprint given its ID.
    """
    url = f"{JIRA_URL}/rest/agile/1.0/sprint/{sprint_id}"
    resp = requests.get(url, auth=(JIRA_EMAIL, JIRA_API_TOKEN))
    resp.raise_for_status()
    data = resp.json()
    return data.get("name", f"Sprint_{sprint_id}")

def get_sprint_dates(sprint_id):
    """
    Fetch the start and end date of a sprint given its ID.
    Returns (start_date, end_date) as strings (YYYY-MM-DD) or (None, None) if not available.
    """
    url = f"{JIRA_URL}/rest/agile/1.0/sprint/{sprint_id}"
    resp = requests.get(url, auth=(JIRA_EMAIL, JIRA_API_TOKEN))
    resp.raise_for_status()
    data = resp.json()
    start = data.get("startDate")
    end = data.get("endDate")
    # Format as YYYY-MM-DD if present
    def fmt(dt):
        if not dt:
            return None
        return dt[:10]
    return (fmt(start), fmt(end))

def group_issues_by_label(issues):
    """
    Group issues by the first matching label from a predefined list (case-insensitive).
    Only includes issues of type 'story' or 'task'.
    Returns a dict: label -> list of issues, in the order of the label list, with 'Other' for unmatched.
    """
    LABEL_MAP = {
        "nlms": "Dutch Platform(s)",
        "iems": "Irish Platform(s)",
        "esms": "Spanish Platform(s)",
        "ukms": "UK Platform(s)",
        "s&a-mpc": "MPC",
        "s&a_mgt": "Management tasks",
        "fims": "Finnish Platform(s)",
    }
    LABEL_ORDER = list(LABEL_MAP.keys())
    grouped = {LABEL_MAP[label]: [] for label in LABEL_ORDER}
    grouped["Other"] = []
    for issue in issues:
        fields = issue["fields"]
        issuetype = fields["issuetype"]["name"].lower()
        if issuetype not in ["story", "task"]:
            continue
        labels = [l.lower() for l in fields.get("labels", [])]
        found = False
        for label_key in LABEL_ORDER:
            if label_key in labels:
                grouped[LABEL_MAP[label_key]].append(issue)
                found = True
                break
        if not found:
            grouped["Other"].append(issue)
    # Remove empty groups for cleaner output
    return {k: v for k, v in grouped.items() if v}

def create_presentation(grouped_issues, filename="Sprint_Review.pptx"):
    # Emoji spinner for progress indicator
    emojis = ["⏳", "🕐", "🕑", "🕒", "🕓", "🕔", "⌛", "🕐", "🕑", "🕒", "🕓", "🕔"]
    def show_progress(step, total, prefix=""):
        idx = (step % len(emojis))
        sys.stdout.write(f"\r{prefix}{emojis[idx]} {step}/{total}")
        sys.stdout.flush()
    """
    Create a PowerPoint presentation from grouped Jira issues.
    - Uses a template for consistent layout.
    - Each label gets a slide with issues listed in the BODY placeholder.
    - Adds a summary slide and an upcoming sprint slide, using the correct placeholders.
    - Adds a final 'thanks' slide if available in the template.
    """
    # Use the sprint-template.pptx as the base for the presentation
    # Load the PowerPoint template from the same directory as the script
    template_path = os.path.join(os.path.dirname(__file__), "sprint-template.pptx")
    if not os.path.exists(template_path):
        raise FileNotFoundError(
            "PowerPoint template 'sprint-template.pptx' not found in the same directory as the script. "
            "Please add the template file and try again."
        )
    prs = Presentation(template_path)

    # Helper to get layout by name (case-insensitive)
    def get_layout_by_name(prs, name):
        """
        Helper to get a slide layout by name (case-insensitive).
        Returns the first layout if not found.
        """
        for layout in prs.slide_layouts:
            if layout.name.strip().lower() == name.strip().lower():
                return layout
        return prs.slide_layouts[0]  # fallback

    # Get the main layouts used for slides
    title_slide_layout = get_layout_by_name(prs, "Title Slide")
    title_content_layout = get_layout_by_name(prs, "Title and Content")
    summary_layout = get_layout_by_name(prs, "Title and Content Blue Hexagon")

    # Add a title slide with sprint name and dates
    # Always use the active sprint for the title (from get_current_sprint_id)


    # Add a title slide with sprint name and dates
    # Always use the active sprint for the title (from get_current_sprint_id)
    try:
        active_sprint_id = get_current_sprint_id()
        sprint_name = get_sprint_name(active_sprint_id)
        sprint_start, sprint_end = get_sprint_dates(active_sprint_id)
    except Exception:
        sprint_name = os.path.splitext(os.path.basename(filename))[0]
        sprint_start = None
        sprint_end = None
        pass

    # Unindented: rest of the function
    title_slide = prs.slides.add_slide(title_slide_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = sprint_name
    sprint_dates = None
    if sprint_start and sprint_end:
        sprint_dates = f"{sprint_start} to {sprint_end}"
    elif sprint_start:
        sprint_dates = f"Start: {sprint_start}"
    elif sprint_end:
        sprint_dates = f"End: {sprint_end}"
    if sprint_dates:
        try:
            title_slide.placeholders[1].text = sprint_dates
        except Exception:
            pass
    # Calculate totals for summary slide
    planned_points = 0
    planned_time = 0
    achieved_points = 0
    achieved_time = 0
    for label, issues in grouped_issues.items():
        for issue in issues:
            fields = issue["fields"]
            story_points = fields.get("customfield_10024")
            # Planned: all stories
            if story_points not in (None, "?") and str(story_points).strip() != "":
                try:
                    planned_points += float(story_points)
                except Exception:
                    pass
            # Planned time: prefer original estimate
            time_estimate = None
            if fields.get("timeoriginalestimate") not in (None, "", "?"):
                time_estimate = fields.get("timeoriginalestimate")
            elif fields.get("timetracking") and isinstance(fields["timetracking"], dict):
                time_estimate = fields["timetracking"].get("originalEstimateSeconds")
            if time_estimate not in (None, "", "?"):
                try:
                    planned_time += int(time_estimate)
                except Exception:
                    pass
            # Achieved: only if status is done
            status = fields.get("status", {}).get("name", "").lower()
            if status in ("done", "closed", "resolved"):
                if story_points not in (None, "?") and str(story_points).strip() != "":
                    try:
                        achieved_points += float(story_points)
                    except Exception:
                        pass
                        pass
                # Achieved time: prefer logged time
                time_logged = None
                if fields.get("timetracking") and isinstance(fields["timetracking"], dict):
                    time_logged = fields["timetracking"].get("timeSpentSeconds")
                if time_logged not in (None, "", "?"):
                    try:
                        achieved_time += int(time_logged)
                    except Exception:
                        pass

    # Create a slide for each label
    total_labels = len(grouped_issues)
    for i, (label, issues) in enumerate(grouped_issues.items(), 1):
        show_progress(i, total_labels, prefix="Slides: ")
        slide = prs.slides.add_slide(title_content_layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"{label}"
        status_map = {
            'done': '✔️',
            'notcompleted': '🛠️',
            'in progress': '🛠️',
            'inprogress': '🛠️',
            'punted': '🧹',
            'added': '➕',
            'added after sprint start': '➕',
            'cancelled': '❌',
            'canceled': '❌',
        }
        issue_lines = []
        # Get the sprint start date for added-after-sprint detection
        sprint_start_dt = None
        active_sprint_id = None
        try:
            active_sprint_id = get_current_sprint_id()
            sprint_start, _ = get_sprint_dates(active_sprint_id)
            if sprint_start:
                from datetime import datetime
                sprint_start_dt = datetime.fromisoformat(sprint_start)
        except Exception:
            pass
        for idx, issue in enumerate(issues):
            key = issue["key"]
            fields = issue["fields"]
            summary = fields.get("summary", "")
            status_name = fields.get("status", {}).get("name", "").lower()
            # Detect if issue was added to the sprint after the sprint started
            added_after_sprint = False
            if sprint_start_dt and active_sprint_id:
                # Check if the sprint was added to the issue after the sprint started
                changelog = issue.get('changelog')
                found_added = False
                if changelog and 'histories' in changelog:
                    from datetime import datetime, timezone
                    for history in changelog['histories']:
                        for item in history.get('items', []):
                            if item.get('field') == 'Sprint':
                                # Check if this sprint was added (look for id in to/from)
                                to_sprint = item.get('to')
                                from_sprint = item.get('from')
                                # 'to' and 'from' can be comma-separated lists of sprint IDs
                                to_ids = [s.strip() for s in str(to_sprint).split(',')] if to_sprint else []
                                from_ids = [s.strip() for s in str(from_sprint).split(',')] if from_sprint else []
                                if str(active_sprint_id) in to_ids and str(active_sprint_id) not in from_ids:
                                    # When was it added?
                                    added_dt = None
                                    try:
                                        added_dt = datetime.fromisoformat(history['created'].replace('Z', '+00:00'))
                                    except Exception:
                                        pass
                                    if added_dt and added_dt > sprint_start_dt:
                                        added_after_sprint = True
                                        found_added = True
                                        break
                        if found_added:
                            break
                # If no changelog, fallback to created date (less accurate)
                elif fields.get('created'):
                    from datetime import datetime, timezone
                    try:
                        created_dt = datetime.fromisoformat(fields['created'].replace('Z', '+00:00'))
                        if created_dt > sprint_start_dt:
                            added_after_sprint = True
                    except Exception:
                        pass
            # Normalize spelling for cancelled/canceled, handle in progress, and added after sprint start
            if status_name in ("cancelled", "canceled"):
                icon = status_map.get("cancelled", '')
            elif status_name.replace(" ","") == "inprogress":
                icon = status_map.get("in progress", '')
            elif added_after_sprint:
                icon = status_map.get("added after sprint start", status_map.get("added", ''))
            else:
                icon = status_map.get(status_name, '')
            assignee = fields.get("assignee")
            display_name = ""
            if assignee and isinstance(assignee, dict):
                display_name = assignee.get("displayName", "")
            story_points = fields.get("customfield_10024")
            time_value = None
            time_label = None
            if fields.get("timetracking") and isinstance(fields["timetracking"], dict):
                time_logged = fields["timetracking"].get("timeSpentSeconds")
                if time_logged not in (None, "", "?"):
                    time_value = time_logged
                    time_label = "Time Logged"
            if time_value is None:
                if fields.get("timeoriginalestimate") not in (None, "", "?"):
                    time_value = fields.get("timeoriginalestimate")
                    time_label = "Time Estimate"
                elif fields.get("timetracking") and isinstance(fields["timetracking"], dict):
                    time_estimate = fields["timetracking"].get("originalEstimateSeconds")
                    if time_estimate not in (None, "", "?"):
                        time_value = time_estimate
                        time_label = "Time Estimate"
            details = []
            if time_value not in (None, "?") and str(time_value).strip() != "":
                try:
                    hours = int(time_value) // 3600
                    minutes = (int(time_value) % 3600) // 60
                    details.append(f"T: {hours}h {minutes}m")
                except Exception:
                    details.append(f"T: {time_value} seconds")
            elif story_points not in (None, "?") and str(story_points).strip() != "":
                details.append(f"P: {story_points}")
            issue_text = f"{key} {icon}: {summary}"
            if display_name:
                issue_text += f" {display_name}"
            if details:
                issue_text += f" ({', '.join(details)})"
            issue_lines.append(issue_text)
        try:
            slide.placeholders[15].text = '\n'.join(issue_lines)
            for paragraph in slide.placeholders[15].text_frame.paragraphs:
                paragraph.font.size = Pt(18)
        except (KeyError, IndexError, AttributeError):
            pass

    # Add summary slide
    show_progress(total_labels + 1, total_labels + 3, prefix="Slides: ")
    summary_slide = prs.slides.add_slide(summary_layout)
    if summary_slide.shapes.title:
        summary_slide.shapes.title.text = "Sprint Summary"
    def format_time(seconds):
        """Format seconds as hours and minutes string."""
        hours = int(seconds) // 3600
        minutes = (int(seconds) % 3600) // 60
        return f"{hours}h {minutes}m"

    summary_text = (
        f"Planned:\n"
        f"  P: {planned_points}\n"
        f"  T: {format_time(planned_time)}\n\n"
        f"Achieved:\n"
        f"  P: {achieved_points}\n"
        f"  T: {format_time(achieved_time)}"
    )
    # Place summary in the correct placeholder (idx 14)
    try:
        summary_slide.placeholders[14].text = summary_text
        for paragraph in summary_slide.placeholders[14].text_frame.paragraphs:
            paragraph.font.size = Pt(18)
    except (KeyError, IndexError, AttributeError):
        pass

    # Add upcoming sprint slide
    show_progress(total_labels + 2, total_labels + 3, prefix="Slides: ")
    # (same logic as label slides, but for upcoming issues)
    upcoming_sprint_id = get_upcoming_sprint_id()
    if upcoming_sprint_id:
        upcoming_issues = get_issues(upcoming_sprint_id)
        upcoming_grouped = group_issues_by_label(upcoming_issues)
        all_upcoming_issues = []
        for issues in upcoming_grouped.values():
            all_upcoming_issues.extend(issues)
        if all_upcoming_issues:
            slide = prs.slides.add_slide(title_content_layout)
            if slide.shapes.title:
                slide.shapes.title.text = "Upcoming"
            # Compose all issue lines for upcoming
            issue_lines = []
            for idx, issue in enumerate(all_upcoming_issues):
                key = issue["key"]
                fields = issue["fields"]
                summary = fields.get("summary", "")
                assignee = fields.get("assignee")
                display_name = ""
                if assignee and isinstance(assignee, dict):
                    display_name = assignee.get("displayName", "")
                story_points = fields.get("customfield_10024")
                time_value = None
                time_label = None
                if fields.get("timetracking") and isinstance(fields["timetracking"], dict):
                    time_logged = fields["timetracking"].get("timeSpentSeconds")
                    if time_logged not in (None, "", "?"):
                        time_value = time_logged
                        time_label = "Time Logged"
                if time_value is None:
                    if fields.get("timeoriginalestimate") not in (None, "", "?"):
                        time_value = fields.get("timeoriginalestimate")
                        time_label = "Time Estimate"
                    elif fields.get("timetracking") and isinstance(fields["timetracking"], dict):
                        time_estimate = fields["timetracking"].get("originalEstimateSeconds")
                        if time_estimate not in (None, "", "?"):
                            time_value = time_estimate
                            time_label = "Time Estimate"
                details = []
                if time_value not in (None, "?") and str(time_value).strip() != "":
                    try:
                        hours = int(time_value) // 3600
                        minutes = (int(time_value) % 3600) // 60
                        details.append(f"T: {hours}h {minutes}m")
                    except Exception:
                        details.append(f"T: {time_value} seconds")
                elif story_points not in (None, "?") and str(story_points).strip() != "":
                    details.append(f"P: {story_points}")
                issue_text = f"{key}: {summary}"
                if display_name:
                    issue_text += f" {display_name}"
                if details:
                    issue_text += f" ({', '.join(details)})"
                issue_lines.append(issue_text)
            # Place all issue lines in the main content placeholder by index
            try:
                slide.placeholders[15].text = '\n'.join(issue_lines)
                for paragraph in slide.placeholders[15].text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
            except (KeyError, IndexError, AttributeError):
                pass

    # Add a final empty "thanks" slide if available
    show_progress(total_labels + 3, total_labels + 3, prefix="Slides: ")
    thanks_layout = get_layout_by_name(prs, "thanks")
    prs.slides.add_slide(thanks_layout)

    # Save the presentation
    import time as _time
    import glob
    import os as _os
    def is_pptx_locked(filename):
        # Check for a lock file in the same directory (e.g., '~$filename')
        dir_name = _os.path.dirname(filename) or '.'
        base_name = _os.path.basename(filename)
        lock_pattern = _os.path.join(dir_name, '~$' + base_name)
        return _os.path.exists(lock_pattern)

    def try_save_pptx(prs, filename):
        while True:
            if is_pptx_locked(filename):
                print(f"\nWARNING: A lock file for '{filename}' was found (likely open in PowerPoint). Please close the file and press Enter to continue...")
                input()
                _time.sleep(1)
            try:
                prs.save(filename)
                break
            except PermissionError:
                print(f"\nERROR: The file '{filename}' is open in PowerPoint or another program. Please close it and press Enter to continue...")
                input()
                _time.sleep(1)
    try_save_pptx(prs, filename)
    sys.stdout.write("\r" + " " * 40 + "\r")  # Clear progress line
    print(f"Presentation saved as {filename}")
    # (No avatar cleanup needed)

if __name__ == "__main__":
    # Emoji spinner for progress indicator
    import threading
    import itertools
    import time

    spinner_running = True
    spinner_message = ["Starting..."]
    def spinner_func():
        emojis = ["⏳", "🕐", "🕑", "🕒", "🕓", "🕔", "⌛", "🕐", "🕑", "🕒", "🕓", "🕔"]
        for idx in itertools.cycle(range(len(emojis))):
            if not spinner_running:
                break
            sys.stdout.write(f"\r{emojis[idx]} {spinner_message[0]}   ")
            sys.stdout.flush()
            time.sleep(0.15)
        sys.stdout.write("\r" + " " * 40 + "\r")

    spinner_thread = threading.Thread(target=spinner_func)
    spinner_thread.start()

    try:
        spinner_message[0] = "Fetching current sprint ID..."
        sprint_id = get_current_sprint_id()
        spinner_message[0] = "Fetching sprint name..."
        sprint_name = get_sprint_name(sprint_id)
        spinner_message[0] = "Fetching issues for current sprint..."
        issues = get_issues(sprint_id)
        spinner_message[0] = "Grouping issues by label..."
        grouped = group_issues_by_label(issues)
        spinner_message[0] = "Creating PowerPoint presentation..."
        # Sanitize sprint_name for filename (remove problematic chars)
        import re
        safe_sprint_name = re.sub(r'[^\w\-_\. ]', '_', sprint_name)
        filename = f"{safe_sprint_name}.pptx"
        create_presentation(grouped, filename=filename)
        spinner_message[0] = "Done!"
    finally:
        spinner_running = False
        spinner_thread.join()
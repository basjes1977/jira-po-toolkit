"""
jpt_presentation.py
------------------
Contains PowerPoint generation logic for Jira sprints.
Split from jpt.py for modularity.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import logging

logger = logging.getLogger("jpt_presentation")

def create_presentation(
    grouped_issues,
    sprint_name,
    sprint_start,
    sprint_end,
    filename="Sprint_Review.pptx",
    epic_map=None,
    epic_goals=None,
    planned_items=None,
    velocity_history=None,
):
    """
    Create a PowerPoint presentation from grouped Jira issues.
    - Each label gets a slide with issues listed in the BODY placeholder.
    - Adds a summary slide showing accomplished work (no points).
    - Adds a final 'thanks' slide if available in the template.
    """
    template_path = os.path.join(os.path.dirname(__file__), "sprint-template.pptx")
    if not os.path.exists(template_path):
        raise FileNotFoundError(
            "PowerPoint template 'sprint-template.pptx' not found in the same directory as the script. "
            "Please add the template file and try again."
        )
    prs = Presentation(template_path)

    def get_layout_by_name(prs, name):
        for layout in prs.slide_layouts:
            if layout.name.strip().lower() == name.strip().lower():
                return layout
        return prs.slide_layouts[0]

    # Helper to extract an epic identifier/name from an issue's fields.
    def detect_epic_name(fields):
        # Common epic link field keys include: customfield_10008, epic, Epic Link, etc.
        # Try known keys first then fall back to scanning field names that contain 'epic'.
        candidates = [
            "customfield_10008",
            "customfield_10006",
            "epic",
            "Epic Link",
            "epic_link",
        ]
        for key in candidates:
            if key in fields and fields.get(key):
                val = fields.get(key)
                # If it's a dict, try name or key
                if isinstance(val, dict):
                    return val.get("key") or val.get("name") or str(val)
                # If it's a list, pick first element
                if isinstance(val, list) and val:
                    first = val[0]
                    if isinstance(first, dict):
                        return first.get("key") or first.get("name") or str(first)
                    return str(first)
                return str(val)
        # Scan all fields for anything with 'epic' in the key
        for k, v in fields.items():
            if "epic" in k.lower() and v:
                if isinstance(v, dict):
                    return v.get("key") or v.get("name") or str(v)
                return str(v)
        return None

    title_slide_layout = get_layout_by_name(prs, "Title Slide")
    title_content_layout = get_layout_by_name(prs, "Title and Content")
    summary_layout = get_layout_by_name(prs, "Title and Content Blue Hexagon")

    # Title slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = sprint_name
    sprint_dates = None
    if sprint_start and sprint_end:
        sprint_dates = f"{sprint_start} to {sprint_end}"
    elif sprint_start:
        sprint_dates = f"Start: {sprint_start}"
    elif sprint_end:
        sprint_dates = f"End: {sprint_end}"
    if sprint_dates:
        try:
            title_slide.placeholders[1].text = sprint_dates
        except Exception:
            pass

    # Slides for each label
    def _paginate_list(lst, n):
        for i in range(0, len(lst), n):
            yield lst[i:i+n]

    for label, issues in grouped_issues.items():
        # Build the textual lines for the label slides
        issue_lines = []
        for issue in issues:
            key = issue.get("key")
            fields = issue.get("fields", {})
            summary = fields.get("summary", "")
            status_name = fields.get("status", {}).get("name", "")
            assignee = fields.get("assignee")
            display_name = ""
            if assignee and isinstance(assignee, dict):
                display_name = assignee.get("displayName", "")
            # Check if issue was added mid-sprint
            mid_sprint_marker = "➕ " if issue.get("_added_mid_sprint") else ""
            issue_text = f"{mid_sprint_marker}{key}: {summary}"
            if display_name:
                issue_text += f" {display_name}"
            if status_name:
                issue_text += f" [{status_name}]"
            issue_lines.append(issue_text)

        # Paginate: max 5 stories per slide for label slides
        pages = list(_paginate_list(issue_lines, 5)) if issue_lines else [[]]
        total = len(pages)
        for idx, page_lines in enumerate(pages, start=1):
            slide = prs.slides.add_slide(title_content_layout)
            title_text = f"{label}"
            if total > 1:
                title_text = f"{label} ({idx}/{total})"
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            try:
                # Try placing into common placeholder indexes; set story font size to 16pt
                placed = False
                for ph in (15, 14, 1, 2, 0):
                    try:
                        slide.placeholders[ph].text = '\n'.join(page_lines)
                        for paragraph in slide.placeholders[ph].text_frame.paragraphs:
                            paragraph.font.size = Pt(16)
                        placed = True
                        break
                    except Exception:
                        continue
                if not placed and page_lines:
                    # Fallback: add a textbox
                    from pptx.util import Inches
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.0))
                    tf = txBox.text_frame
                    tf.text = '\n'.join(page_lines)
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(16)
            except (KeyError, IndexError, AttributeError):
                pass

    # Create one slide per epic: show epic title (from epic_map), Goal (if available), and list stories/tasks
    epic_items = {}
    for label, issues in grouped_issues.items():
        for issue in issues:
            fields = issue.get("fields", {})
            epic_name = detect_epic_name(fields)
            # Only include issues that have an epic defined; skip those without one
            if not epic_name or epic_name == "None":
                continue
            epic_items.setdefault(epic_name, []).append(issue)

    if not epic_items:
        # fallback single slide when no epics found
        summary_slide = prs.slides.add_slide(summary_layout)
        if summary_slide.shapes.title:
            summary_slide.shapes.title.text = "Sprint Epic Progress"
        try:
            summary_slide.placeholders[14].text = "No stories or tasks found for this sprint."
        except Exception:
            pass
    else:
        # Helpers
        def chunk_list(lst, n):
            for i in range(0, len(lst), n):
                yield lst[i:i+n]

        def truncate(text, length=100):
            if not text:
                return ""
            s = str(text).strip()
            if len(s) <= length:
                return s
            return s[:length-1].rsplit(' ', 1)[0] + '…'

        from pptx.util import Inches

        # Build structured data: initiative -> list of (epic_display, items_sorted)
        initiatives = {}
        # We'll also keep a mapping from initiative_display -> description (if available)
        initiative_descriptions = {}
        for epic, items in epic_items.items():
            display_epic = epic_map.get(epic, epic) if epic_map else epic
            init_text = None
            init_desc = None
            if epic_goals and epic in epic_goals:
                val = epic_goals[epic]
                if isinstance(val, dict):
                    init_text = val.get("display")
                    init_desc = val.get("description")
                else:
                    init_text = str(val)
            initiative_display = init_text or "No Initiative"
            if initiative_display not in initiative_descriptions and init_desc:
                initiative_descriptions[initiative_display] = init_desc
            items_sorted = sorted(items, key=lambda it: it.get("key") or "")
            initiatives.setdefault(initiative_display, []).append((display_epic, items_sorted))

        # -----------------------------
        # Section B: Per-initiative slides (user's preference)
        # -----------------------------
        # One slide per initiative, repeating the initiative header on split slides when necessary.
        for initiative_display, epics in initiatives.items():
            # Flatten epics into chunks of up to 5 epics per slide to avoid overcrowding
            for ep_batch in chunk_list(epics, 5):
                slide = prs.slides.add_slide(title_content_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = f"Initiative: {initiative_display}"
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.0))
                tf = txBox.text_frame
                tf.clear()
                # If we have a description for this initiative, render the first paragraph (truncated)
                idesc = initiative_descriptions.get(initiative_display)
                if idesc:
                    desc_text = truncate(idesc, 300)
                    p_desc = tf.add_paragraph()
                    p_desc.text = desc_text
                    p_desc.font.size = Pt(12)
                    p_desc.font.italic = True
                    tf.add_paragraph()
                for display_epic, items in ep_batch:
                    p = tf.add_paragraph()
                    p.text = display_epic
                    p.font.size = Pt(16)
                    p.font.bold = True
                    # stories
                    for issue in items:
                        f = issue.get("fields", {})
                        status = f.get("status", {}).get("name", "").lower()
                        done = status in ("done", "closed", "resolved")
                        key = issue.get("key")
                        summary = truncate(f.get("summary", ""), 120)
                        mark = "✔️" if done else "—"
                        # Check if issue was added mid-sprint
                        mid_sprint_marker = " ➕" if issue.get("_added_mid_sprint") else ""
                        s = tf.add_paragraph()
                        s.text = f"{mark}{mid_sprint_marker} {key}: {summary} [{status}]"
                        s.level = 1
                        s.font.size = Pt(11)
                    prog = tf.add_paragraph()
                    prog.text = f"Progress: {sum(1 for it in items if it.get('fields', {}).get('status', {}).get('name','').lower() in ('done','closed','resolved'))}/{len(items)} done"
                    prog.level = 1
                    prog.font.size = Pt(11)
                    tf.add_paragraph()

    # Thanks slide
    # --- Planned items (next sprint + in-progress) ---
    if planned_items:
        # Build textual lines for planned items (only show KEY: summary)
        planned_lines = []
        for issue in planned_items:
            key = issue.get("key")
            fields = issue.get("fields", {})
            summary = fields.get("summary", "")
            issue_text = f"{key}: {summary}"
            planned_lines.append(issue_text)

        # Paginate planned items: 5 per slide
        pages = list(_paginate_list(planned_lines, 5)) if planned_lines else []
        total = len(pages)
        for idx, page_lines in enumerate(pages, start=1):
            slide = prs.slides.add_slide(title_content_layout)
            title_text = "Planned for next sprint"
            if total > 1:
                title_text = f"Planned for next sprint ({idx}/{total})"
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            try:
                placed = False
                for ph in (15, 14, 1, 2, 0):
                    try:
                        slide.placeholders[ph].text = '\n'.join(page_lines)
                        for paragraph in slide.placeholders[ph].text_frame.paragraphs:
                            paragraph.font.size = Pt(16)
                        placed = True
                        break
                    except Exception:
                        continue
                if not placed and page_lines:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.0))
                    tf = txBox.text_frame
                    tf.text = '\n'.join(page_lines)
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(16)
            except (KeyError, IndexError, AttributeError):
                pass

    # Velocity slide (based on recent sprint history)
    if velocity_history:
        slide = prs.slides.add_slide(title_content_layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"Velocity (Last {len(velocity_history)} Sprints)"
        history = list(reversed(velocity_history))  # oldest -> newest
        categories = [entry.get("name") or "Sprint" for entry in history]
        points_series = [(entry.get("points") or 0.0) for entry in history]
        hours_series = [round(((entry.get("time_seconds") or 0) / 3600.0), 1) for entry in history]

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("Completed Points", points_series)
        chart_data.add_series("Logged Hours", hours_series)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(0.5),
            Inches(1.4),
            Inches(9),
            Inches(4.3),
            chart_data,
        ).chart
        chart.has_legend = True
        if chart.category_axis:
            chart.category_axis.tick_labels.font.size = Pt(11)
        if chart.value_axis:
            chart.value_axis.tick_labels.font.size = Pt(11)

        avg_points = sum(points_series) / len(points_series) if points_series else None
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(9), Inches(1.2))
        tf = textbox.text_frame
        tf.text = f"Average completed points: {avg_points:.1f}" if avg_points is not None else "Average completed points: n/a"

    thanks_layout = get_layout_by_name(prs, "thanks")
    prs.slides.add_slide(thanks_layout)

    prs.save(filename)
    print(f"Presentation saved as {filename}")
